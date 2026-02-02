"""
Script pour generer la presentation PowerPoint Rakuten Soutenance
avec le meme contenu que la version HTML.
"""

import subprocess
import sys

# Install python-pptx if not available
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Cm
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.oxml.ns import nsmap
except ImportError:
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
    from pptx import Presentation
    from pptx.util import Inches, Pt, Cm
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.oxml.ns import nsmap

# Couleurs Rakuten
RAKUTEN_RED = RGBColor(0xBF, 0x00, 0x00)
DARK_TEXT = RGBColor(0x1E, 0x29, 0x3B)
MUTED_TEXT = RGBColor(0x64, 0x74, 0x8B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF8, 0xF9, 0xFA)
SUCCESS_GREEN = RGBColor(0x10, 0xB9, 0x81)
WARNING_ORANGE = RGBColor(0xF5, 0x9E, 0x0B)
ACCENT_BLUE = RGBColor(0x25, 0x63, 0xEB)
ACCENT_PURPLE = RGBColor(0x7C, 0x3A, 0xED)


def set_shape_fill(shape, color):
    """Set solid fill color for a shape."""
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def add_title_slide(prs, title, subtitle=""):
    """Add a title slide."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Background
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    set_shape_fill(background, WHITE)
    background.line.fill.background()

    # Red accent line at top
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Pt(6))
    set_shape_fill(line, RAKUTEN_RED)
    line.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Cm(2), Cm(6), Cm(21), Cm(3))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = DARK_TEXT
    p.alignment = PP_ALIGN.CENTER

    if subtitle:
        p = tf.add_paragraph()
        p.text = subtitle
        p.font.size = Pt(20)
        p.font.color.rgb = MUTED_TEXT
        p.alignment = PP_ALIGN.CENTER
        p.space_before = Pt(20)

    return slide


def add_content_slide(prs, badge_text, title, badge_color=RAKUTEN_RED):
    """Add a content slide with header."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Background
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    set_shape_fill(background, WHITE)
    background.line.fill.background()

    # Badge
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(1.5), Cm(1), Cm(3), Cm(0.8))
    set_shape_fill(badge, badge_color)
    badge.line.fill.background()
    badge.text_frame.paragraphs[0].text = badge_text
    badge.text_frame.paragraphs[0].font.size = Pt(10)
    badge.text_frame.paragraphs[0].font.bold = True
    badge.text_frame.paragraphs[0].font.color.rgb = WHITE
    badge.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Title
    title_box = slide.shapes.add_textbox(Cm(5), Cm(0.85), Cm(18), Cm(1.2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = DARK_TEXT

    # Separator line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(1.5), Cm(2.1), Cm(22), Pt(2))
    set_shape_fill(line, RGBColor(0xE2, 0xE8, 0xF0))
    line.line.fill.background()

    return slide


def add_text_box(slide, left, top, width, height, text, font_size=14, bold=False, color=DARK_TEXT, align=PP_ALIGN.LEFT):
    """Add a text box to a slide."""
    box = slide.shapes.add_textbox(Cm(left), Cm(top), Cm(width), Cm(height))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return box


def add_metric_box(slide, left, top, value, label, color=RAKUTEN_RED):
    """Add a metric box."""
    # Box background
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(left), Cm(top), Cm(4.5), Cm(2.5))
    set_shape_fill(box, LIGHT_BG)
    box.line.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)

    # Value
    val_box = slide.shapes.add_textbox(Cm(left), Cm(top + 0.3), Cm(4.5), Cm(1.2))
    tf = val_box.text_frame
    p = tf.paragraphs[0]
    p.text = value
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = color
    p.alignment = PP_ALIGN.CENTER

    # Label
    lbl_box = slide.shapes.add_textbox(Cm(left), Cm(top + 1.5), Cm(4.5), Cm(0.8))
    tf = lbl_box.text_frame
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(10)
    p.font.color.rgb = MUTED_TEXT
    p.alignment = PP_ALIGN.CENTER


def add_table(slide, left, top, width, rows_data, col_widths=None):
    """Add a table to slide."""
    rows = len(rows_data)
    cols = len(rows_data[0])

    table = slide.shapes.add_table(rows, cols, Cm(left), Cm(top), Cm(width), Cm(rows * 0.8)).table

    for i, row_data in enumerate(rows_data):
        for j, cell_text in enumerate(row_data):
            cell = table.cell(i, j)
            cell.text = str(cell_text)

            # Header row
            if i == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RAKUTEN_RED
                cell.text_frame.paragraphs[0].font.color.rgb = WHITE
                cell.text_frame.paragraphs[0].font.bold = True
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE
                cell.text_frame.paragraphs[0].font.color.rgb = DARK_TEXT

            cell.text_frame.paragraphs[0].font.size = Pt(11)

    return table


def create_presentation():
    """Create the full presentation."""
    prs = Presentation()
    prs.slide_width = Cm(25.4)  # 16:9
    prs.slide_height = Cm(14.29)

    # ========== SLIDE 0: COVER ==========
    slide = add_title_slide(prs, "Classification Multimodale", "Produits E-Commerce - Challenge Rakuten France")

    # Metrics row
    add_metric_box(slide, 3, 9.5, "84 916", "Produits", RAKUTEN_RED)
    add_metric_box(slide, 8, 9.5, "27", "Categories", RAKUTEN_RED)
    add_metric_box(slide, 13, 9.5, "92%", "Accuracy Image", SUCCESS_GREEN)
    add_metric_box(slide, 18, 9.5, "83%", "Accuracy Texte", WARNING_ORANGE)

    # Team
    add_text_box(slide, 1, 12.5, 23, 1,
                 "Johan Frachon  -  Liviu Andronic  -  Hery M. Ralaimanantsoa  -  Oussama Akir",
                 font_size=11, color=MUTED_TEXT, align=PP_ALIGN.CENTER)
    add_text_box(slide, 1, 13.2, 23, 0.8, "Mentor : Antoine - Fevrier 2025",
                 font_size=10, color=MUTED_TEXT, align=PP_ALIGN.CENTER)

    # ========== SLIDE 1: SOMMAIRE ==========
    slide = add_content_slide(prs, "NAVIGATION", "Sommaire")

    toc_items = [
        ("01", "Le Challenge Rakuten"),
        ("02", "Desequilibre des Classes"),
        ("03", "Preprocessing (TF-IDF + CNN)"),
        ("04", "Modelisation Texte - 83%"),
        ("05", "Modelisation Image (ML vs DL)"),
        ("06", "Voting System - 92%"),
        ("07", "Fusion Multimodale - 94%"),
        ("08", "Architecture End-to-End"),
        ("09", "Distribution des Classes"),
        ("10", "Exemples de Predictions"),
        ("11", "Conclusion & Perspectives"),
    ]

    for i, (num, title) in enumerate(toc_items):
        col = 0 if i < 6 else 1
        row = i if i < 6 else i - 6
        left = 2 + col * 11
        top = 3 + row * 1.5

        # Number box
        num_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(left), Cm(top), Cm(1.2), Cm(1))
        set_shape_fill(num_box, LIGHT_BG)
        num_box.line.fill.background()
        num_box.text_frame.paragraphs[0].text = num
        num_box.text_frame.paragraphs[0].font.size = Pt(14)
        num_box.text_frame.paragraphs[0].font.bold = True
        num_box.text_frame.paragraphs[0].font.color.rgb = RAKUTEN_RED
        num_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Title
        add_text_box(slide, left + 1.5, top + 0.2, 9, 0.8, title, font_size=13, color=DARK_TEXT)

    # ========== SLIDE 2: CHALLENGE ==========
    slide = add_content_slide(prs, "CONTEXTE", "Le Challenge Rakuten France")

    add_text_box(slide, 1.5, 2.8, 10, 2,
                 "Objectif\nClassifier automatiquement des produits e-commerce parmi 27 categories en utilisant le texte et l'image.",
                 font_size=12, color=DARK_TEXT)

    add_text_box(slide, 1.5, 5.5, 10, 3,
                 "Enjeux Business\n- Ameliorer l'experience de recherche\n- Reduire le taux de rebond\n- Automatiser millions de produits/jour\n- Traiter donnees multilingues",
                 font_size=11, color=DARK_TEXT)

    add_metric_box(slide, 13, 3, "84 916", "Produits Train", RAKUTEN_RED)
    add_metric_box(slide, 18, 3, "13 812", "Produits Test", RAKUTEN_RED)
    add_metric_box(slide, 13, 6, "27", "Categories", SUCCESS_GREEN)
    add_metric_box(slide, 18, 6, "35%", "Descriptions NaN", WARNING_ORANGE)

    add_text_box(slide, 13, 9.5, 10, 1.5,
                 "Metrique: F1-Score weighted\nPrend en compte le desequilibre des classes.",
                 font_size=11, color=ACCENT_BLUE)

    # ========== SLIDE 3: DESEQUILIBRE ==========
    slide = add_content_slide(prs, "DONNEES", "Desequilibre des Classes", ACCENT_BLUE)

    add_text_box(slide, 1.5, 2.8, 20, 1,
                 "L'analyse revele un desequilibre significatif constituant un defi majeur du projet.",
                 font_size=12, color=MUTED_TEXT)

    table_data = [
        ["Rang", "Categorie", "Effectif", "%"],
        ["1", "Piscines", "10 217", "12.0%"],
        ["2", "Mobilier", "5 076", "6.0%"],
        ["3", "Deco", "4 996", "5.9%"],
        ["...", "...", "...", "..."],
        ["27", "Consoles", "761", "0.9%"],
    ]
    add_table(slide, 1.5, 4, 10, table_data)

    add_metric_box(slide, 13, 4, "13x", "Ratio Max/Min", WARNING_ORANGE)
    add_metric_box(slide, 18, 4, "3 145", "Moyenne/Classe", RAKUTEN_RED)

    add_text_box(slide, 13, 7.5, 10, 2,
                 "Solutions adoptees:\n- class_weight='balanced'\n- F1-Score weighted\n- Oversampling cible",
                 font_size=11, color=SUCCESS_GREEN)

    # ========== SLIDE 4: PREPROCESSING ==========
    slide = add_content_slide(prs, "PIPELINE", "Preprocessing")

    add_text_box(slide, 1.5, 3, 10, 1, "Pipeline Texte - TF-IDF", font_size=16, bold=True, color=DARK_TEXT)
    add_text_box(slide, 1.5, 4, 10, 1.5,
                 "HTML+NaN -> Clean -> Concat -> TF-IDF\n\nFeatureUnion: Word n-grams (1,2) + Char n-grams (3,5)\n280K features",
                 font_size=11, color=DARK_TEXT)

    add_text_box(slide, 1.5, 7, 10, 1, "Pipeline Image - Transfer Learning", font_size=16, bold=True, color=DARK_TEXT)
    add_text_box(slide, 1.5, 8, 10, 1.5,
                 "500x500 -> Resize 224x224 -> Normalize -> CNN -> Features\n\nResNet50 (2048) / EfficientNet-B0 / ImageNet weights",
                 font_size=11, color=DARK_TEXT)

    add_text_box(slide, 13, 4, 10, 3,
                 "Justification\n\nLa combinaison word + char TF-IDF capture semantique et morphologie, robuste aux fautes d'orthographe et textes multilingues.",
                 font_size=11, color=ACCENT_BLUE)

    # ========== SLIDE 5: TEXTE ==========
    slide = add_content_slide(prs, "TEXTE", "Modelisation Texte - LinearSVC", SUCCESS_GREEN)

    table_data = [
        ["Modele", "F1-Score", "Temps"],
        ["Logistic Regression", "0.79", "~2 min"],
        ["Multinomial NB", "0.76", "~30 sec"],
        ["LinearSVC (C=0.5)", "0.83", "~3 min"],
        ["Random Forest", "0.68", "~15 min"],
    ]
    add_table(slide, 1.5, 3.5, 11, table_data)

    add_metric_box(slide, 15, 4, "83%", "Accuracy LinearSVC", SUCCESS_GREEN)

    add_text_box(slide, 13, 7.5, 10, 2,
                 "Champion Texte: LinearSVC\nclass_weight='balanced', C=0.5\nOptimise par GridSearch\nF1-Score weighted: 0.827",
                 font_size=11, color=SUCCESS_GREEN)

    # ========== SLIDE 6: IMAGE BENCHMARK ==========
    slide = add_content_slide(prs, "IMAGE", "Modelisation Image - ML vs Deep Learning", ACCENT_PURPLE)

    table_data = [
        ["Modele", "F1-Score", "Temps", "Hardware"],
        ["Random Forest", "0.71", "~30 min", "CPU"],
        ["XGBoost GPU", "0.72", "~10 min", "GPU"],
        ["XGBoost Heavy", "0.765", "6 heures", "CPU 128GB"],
        ["MLP (Adam+GELU)", "0.914", "58 sec", "GPU"],
    ]
    add_table(slide, 1.5, 3, 14, table_data)

    add_text_box(slide, 1.5, 8, 10, 2,
                 "Plafond ML: Les modeles classiques plafonnent a ~76% malgre 6h de calcul.",
                 font_size=11, color=WARNING_ORANGE)

    add_text_box(slide, 13, 8, 10, 2,
                 "Rupture Deep Learning: MLP atteint 91.4% en 58 secondes sur GPU RTX 4070.",
                 font_size=11, color=SUCCESS_GREEN)

    # ========== SLIDE 7: VOTING SYSTEM ==========
    slide = add_content_slide(prs, "INNOVATION", "Voting System - Modele Final Image", RAKUTEN_RED)

    add_text_box(slide, 1.5, 2.8, 22, 1,
                 "Exploitation de la complementarite des architectures : les modeles font des erreurs differentes.",
                 font_size=12, color=MUTED_TEXT, align=PP_ALIGN.CENTER)

    # Model boxes
    models = [
        ("DINOv3", "79.1%", "Poids: 0.40"),
        ("XGBoost", "80.1%", "Poids: 0.35"),
        ("EfficientNet", "~75%", "Poids: 0.25"),
    ]

    for i, (name, score, weight) in enumerate(models):
        left = 2 + i * 6
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(left), Cm(4.5), Cm(5), Cm(3))
        set_shape_fill(box, LIGHT_BG)
        box.line.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)

        add_text_box(slide, left + 0.3, 4.8, 4.4, 0.6, name, font_size=14, bold=True, color=DARK_TEXT, align=PP_ALIGN.CENTER)
        add_text_box(slide, left + 0.3, 5.5, 4.4, 0.8, score, font_size=20, bold=True, color=RAKUTEN_RED, align=PP_ALIGN.CENTER)
        add_text_box(slide, left + 0.3, 6.5, 4.4, 0.5, weight, font_size=10, color=MUTED_TEXT, align=PP_ALIGN.CENTER)

    # Final result box
    final_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(20), Cm(4.5), Cm(4), Cm(3))
    set_shape_fill(final_box, RAKUTEN_RED)
    final_box.line.fill.background()

    add_text_box(slide, 20.2, 4.8, 3.6, 0.6, "VOTING", font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide, 20.2, 5.5, 3.6, 0.8, "92%", font_size=24, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide, 20.2, 6.5, 3.6, 0.5, "Final", font_size=10, color=WHITE, align=PP_ALIGN.CENTER)

    add_text_box(slide, 1.5, 8.5, 10, 1.5,
                 "Orthogonalite des Erreurs: Ou DINOv3 se trompe, XGBoost peut avoir raison.",
                 font_size=11, color=ACCENT_BLUE)

    add_text_box(slide, 13, 8.5, 10, 1.5,
                 "Calibration (Sharpening): p^3/Sum(p^3) pour renforcer la confiance.",
                 font_size=11, color=SUCCESS_GREEN)

    # ========== SLIDE 8: FUSION ==========
    slide = add_content_slide(prs, "FUSION", "Fusion Multimodale - Late Fusion", WARNING_ORANGE)

    add_text_box(slide, 1.5, 3.5, 22, 1.5,
                 "Image (92%) x 0.6  +  Texte (83%) x 0.4  =  Final (~94%)",
                 font_size=18, bold=True, color=DARK_TEXT, align=PP_ALIGN.CENTER)

    add_text_box(slide, 1.5, 6, 10, 2.5,
                 "Exemple de Synergie\n\nImage: Forme ronde bleue -> 'Piscine'\nTexte: 'DVD Le Grand Bleu' -> 'DVD'\nFusion: Le texte corrige l'erreur visuelle",
                 font_size=11, color=DARK_TEXT)

    table_data = [
        ["Configuration", "Accuracy"],
        ["Texte seul", "83%"],
        ["Image seule", "92%"],
        ["Fusion Multimodale", "~94%"],
    ]
    add_table(slide, 13, 6, 9, table_data)

    add_text_box(slide, 1.5, 10, 22, 1,
                 "Gain: +2 points par rapport au meilleur modele seul + amelioration de la robustesse.",
                 font_size=11, color=SUCCESS_GREEN, align=PP_ALIGN.CENTER)

    # ========== SLIDE 9: ARCHITECTURE ==========
    slide = add_content_slide(prs, "VUE GLOBALE", "Architecture End-to-End", ACCENT_BLUE)

    add_text_box(slide, 1.5, 3, 22, 6,
                 """
PRODUIT (Texte + Image)
        |
        v
   +----+----+
   |         |
   v         v
TF-IDF    Transfer Learning
280K       CNN Features
   |         |
   v         v
LinearSVC  Voting System
  83%     (DINOv3+XGBoost+EffNet)
   |         92%
   |         |
   +----+----+
        |
        v
   LATE FUSION
   (alpha=0.6 img, 0.4 txt)
        |
        v
   PREDICTION ~94%
""",
                 font_size=10, color=DARK_TEXT, align=PP_ALIGN.CENTER)

    # ========== SLIDE 10: DISTRIBUTION ==========
    slide = add_content_slide(prs, "VISUALISATION", "Distribution des 27 Categories", ACCENT_PURPLE)

    add_text_box(slide, 1.5, 3, 22, 1,
                 "Ratio 13:1 entre la classe majoritaire (Piscines: 10K) et minoritaire (Consoles: 761)",
                 font_size=14, color=DARK_TEXT, align=PP_ALIGN.CENTER)

    add_text_box(slide, 1.5, 5, 10, 3,
                 "Defi Majeur\n\nLa classe majoritaire (Piscines) contient 13x plus d'exemples que la minoritaire (Consoles). Sans traitement, le modele ignorerait les classes rares.",
                 font_size=11, color=WARNING_ORANGE)

    add_text_box(slide, 13, 5, 10, 3,
                 "Strategies Adoptees\n\n- class_weight='balanced'\n- F1-Score weighted\n- Oversampling cible sur classes <1000",
                 font_size=11, color=SUCCESS_GREEN)

    # ========== SLIDE 11: EXEMPLES ==========
    slide = add_content_slide(prs, "DEMO", "Exemples de Predictions", SUCCESS_GREEN)

    add_text_box(slide, 1.5, 2.8, 22, 1,
                 "Cas reels montrant la complementarite texte/image et la robustesse du systeme.",
                 font_size=12, color=MUTED_TEXT, align=PP_ALIGN.CENTER)

    examples = [
        ("Cas Clair", "Texte: 'Harry Potter tome 1'\nImage: Couverture livre", "-> Livres (99%)", SUCCESS_GREEN),
        ("Image Decisive", "Texte: 'Accessoire bleu'\nImage: Bouee piscine", "-> Piscine (92%)", ACCENT_BLUE),
        ("Texte Decisif", "Texte: 'DVD Le Grand Bleu'\nImage: Pochette bleue", "-> DVD/Film (87%)", ACCENT_PURPLE),
    ]

    for i, (title, desc, result, color) in enumerate(examples):
        left = 1.5 + i * 7.5

        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(left), Cm(4), Cm(7), Cm(5))
        set_shape_fill(box, LIGHT_BG)
        box.line.color.rgb = color

        add_text_box(slide, left + 0.3, 4.3, 6.4, 0.6, title, font_size=14, bold=True, color=color, align=PP_ALIGN.CENTER)
        add_text_box(slide, left + 0.3, 5.2, 6.4, 2, desc, font_size=10, color=DARK_TEXT, align=PP_ALIGN.CENTER)
        add_text_box(slide, left + 0.3, 7.5, 6.4, 0.8, result, font_size=12, bold=True, color=color, align=PP_ALIGN.CENTER)

    add_text_box(slide, 1.5, 10, 22, 1,
                 "Synergie Multimodale: Quand l'image est ambigue, le texte tranche. Et vice-versa.",
                 font_size=11, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)

    # ========== SLIDE 12: CONCLUSION ==========
    slide = add_title_slide(prs, "Conclusion & Perspectives", "")

    add_metric_box(slide, 3, 4, "92%", "Image Voting", SUCCESS_GREEN)
    add_metric_box(slide, 10.5, 4, "83%", "Texte SVC", RAKUTEN_RED)
    add_metric_box(slide, 18, 4, "~94%", "Fusion", WARNING_ORANGE)

    add_text_box(slide, 1.5, 7.5, 7, 2,
                 "Methodologie\nPipeline complet avec gestion du desequilibre",
                 font_size=11, color=DARK_TEXT, align=PP_ALIGN.CENTER)

    add_text_box(slide, 9.5, 7.5, 7, 2,
                 "Performance\nVoting System innovant a 92%",
                 font_size=11, color=DARK_TEXT, align=PP_ALIGN.CENTER)

    add_text_box(slide, 17.5, 7.5, 7, 2,
                 "Application\nInterface Streamlit fonctionnelle",
                 font_size=11, color=DARK_TEXT, align=PP_ALIGN.CENTER)

    add_text_box(slide, 1.5, 10.5, 22, 1.5,
                 "Perspectives: OCR (+3-5%) | CamemBERT (+5-8%) | Fine-tuning DINOv3 | Deploiement Cloud",
                 font_size=12, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)

    add_text_box(slide, 1.5, 12.5, 22, 1,
                 "Merci ! - Questions ?",
                 font_size=18, bold=True, color=RAKUTEN_RED, align=PP_ALIGN.CENTER)

    # Save
    output_path = r"D:\datascientest\workspace\OCT25_BMLE_RAKUTEN_WS\repo\OCT25_BMLE_RAKUTEN\reports\PRESENTATION_RAKUTEN_SOUTENANCE.pptx"
    prs.save(output_path)
    print(f"[OK] Presentation PowerPoint generee: {output_path}")
    return output_path


if __name__ == "__main__":
    create_presentation()
